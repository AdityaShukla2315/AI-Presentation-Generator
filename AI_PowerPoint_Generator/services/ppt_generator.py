import os
import tempfile
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from services.images import download_image

def generate_ppt(title: str, slides: list, theme: dict) -> str:
    """Generate PowerPoint presentation"""
    prs = Presentation()
    
    # Set slide size to standard 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    add_title_slide(prs, title, theme)
    
    # Add content slides
    for slide_data in slides:
        if slide_data.get('image_url') or slide_data.get('image_path'):
            add_image_slide(prs, slide_data, theme)
        else:
            add_bullet_slide(prs, slide_data, theme)
    
    # Save to temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    temp_file.close()
    
    return temp_file.name

def add_title_slide(prs, title: str, theme: dict):
    """Add title slide to presentation"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Apply theme formatting
    apply_title_formatting(title_shape, theme)
    
    # Set background color if specified
    if theme.get('background_color'):
        set_slide_background(slide, theme['background_color'])

def add_bullet_slide(prs, slide_data: dict, theme: dict):
    """Add bullet point slide to presentation"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = slide_data['title']
    apply_title_formatting(title_shape, theme)
    
    # Add bullet points
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    bullets = slide_data.get('bullets', [])
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.level = 0
        apply_body_formatting(p, theme)
    
    # Add speaker notes if available
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']
    
    # Set background color if specified
    if theme.get('background_color'):
        set_slide_background(slide, theme['background_color'])

def add_image_slide(prs, slide_data: dict, theme: dict):
    """Add slide with image and optional text"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title if present
    if slide_data['title']:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        apply_title_formatting(title_box, theme)
    
    # Download and add image
    image_path = None
    try:
        if slide_data.get('image_url'):
            image_path = download_image(slide_data['image_url'])
        elif slide_data.get('image_path'):
            image_path = slide_data['image_path']
        
        if image_path and os.path.exists(image_path):
            # Calculate image position and size
            left = Inches(1)
            top = Inches(1.5)
            max_width = Inches(11.333)
            max_height = Inches(5)
            
            slide.shapes.add_picture(image_path, left, top, width=max_width, height=max_height)
            
            # Clean up downloaded image
            if slide_data.get('image_url') and image_path != slide_data.get('image_path'):
                try:
                    os.remove(image_path)
                except:
                    pass
    
    except Exception as e:
        logging.error(f"Error adding image to slide: {e}")
        # Fallback to bullet slide if image fails
        add_bullet_slide(prs, slide_data, theme)
        return
    
    # Add bullet points if any
    bullets = slide_data.get('bullets', [])
    if bullets:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(1))
        text_frame = text_box.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"• {bullet}"
            apply_body_formatting(p, theme)
    
    # Add speaker notes if available
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']
    
    # Set background color if specified
    if theme.get('background_color'):
        set_slide_background(slide, theme['background_color'])

def apply_title_formatting(shape, theme: dict):
    """Apply theme formatting to title text"""
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            font = run.font
            font.name = theme.get('font_name', 'Arial')
            font.size = Pt(44)
            font.bold = True
            if theme.get('title_color'):
                font.color.rgb = RGBColor.from_string(theme['title_color'].lstrip('#'))

def apply_body_formatting(paragraph, theme: dict):
    """Apply theme formatting to body text"""
    for run in paragraph.runs:
        font = run.font
        font.name = theme.get('font_name', 'Arial')
        font.size = Pt(24)
        if theme.get('body_color'):
            font.color.rgb = RGBColor.from_string(theme['body_color'].lstrip('#'))

def set_slide_background(slide, color: str):
    """Set slide background color"""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color.lstrip('#'))
    except Exception as e:
        logging.warning(f"Could not set background color: {e}")
